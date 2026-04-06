"""Generate a 5-slide PPTX for the DiscreteEventSimulator repository."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT   = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
YELLOW   = RGBColor(0xFB, 0xBF, 0x24)
CODE_BG  = RGBColor(0x1E, 0x29, 0x3B)

W  = Inches(13.33)   # widescreen 16:9
H  = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color: RGBColor = BG) -> None:
    from pptx.util import Emu
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha: int = None):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_label(slide, text: str, left, top, width, height,
              size: int = 20, bold: bool = False, color: RGBColor = WHITE,
              align=PP_ALIGN.LEFT, italic: bool = False) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic


def add_code_block(slide, code: str, left, top, width, height, size: int = 14) -> None:
    add_rect(slide, left, top, width, height, CODE_BG)
    txb = slide.shapes.add_textbox(
        left + Inches(0.18), top + Inches(0.15),
        width - Inches(0.36), height - Inches(0.3)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = ACCENT
        run.font.name  = "Courier New"


def divider(slide, top, color: RGBColor = ACCENT) -> None:
    add_rect(slide, Inches(0.5), top, Inches(12.33), Inches(0.04), color)


# ── Slide 1 — Title ───────────────────────────────────────────────────────────

def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    # accent bar left edge
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, ACCENT)

    # big title
    add_label(sl, "Discrete Event Simulator",
              Inches(0.5), Inches(1.8), Inches(12), Inches(1.4),
              size=54, bold=True, color=WHITE)

    # subtitle
    add_label(sl, "Queueing network simulation with live CLI visualization",
              Inches(0.5), Inches(3.35), Inches(10), Inches(0.6),
              size=24, color=ACCENT)

    divider(sl, Inches(4.2))

    # three feature pills
    pills = [
        ("M/M/c Queues",       Inches(0.5)),
        ("Jackson Networks",   Inches(4.0)),
        ("Rich CLI Dashboard", Inches(7.7)),
    ]
    for label, left in pills:
        add_rect(sl, left, Inches(4.5), Inches(3.1), Inches(0.55), RGBColor(0x1E,0x29,0x3B))
        add_label(sl, label, left + Inches(0.15), Inches(4.52), Inches(2.9), Inches(0.5),
                  size=18, color=ACCENT, bold=True)

    add_label(sl, "Python 3  ·  pure stdlib DES engine  ·  rich + networkx",
              Inches(0.5), Inches(5.6), Inches(12), Inches(0.4),
              size=14, color=MUTED, italic=True)


# ── Slide 2 — Architecture ────────────────────────────────────────────────────

def slide_architecture(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, ACCENT)

    add_label(sl, "Architecture", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
              size=36, bold=True, color=WHITE)
    divider(sl, Inches(1.15))

    # layer boxes
    layers = [
        ("examples/",    "User scripts — define networks, call net.run()", ACCENT),
        ("des/network/", "QueueingNetwork — wires nodes, owns DiGraph, routing", GREEN),
        ("des/nodes/",   "Source · MMcServer · Sink — Source uses Poisson arrivals (Exp inter-arrival) by default; pluggable via inter_arrival_fn", YELLOW),
        ("des/engine/",  "Scheduler (heapq) · Simulation clock · Event dispatch", RGBColor(0xF4,0x72,0x18)),
        ("des/stats/",   "Welford + time-weighted accumulators — W, Wq, L, Lq", MUTED),
    ]

    top = Inches(1.4)
    for name, desc, color in layers:
        add_rect(sl, Inches(0.5), top, Inches(12.3), Inches(0.72), CODE_BG)
        add_rect(sl, Inches(0.5), top, Inches(0.08), Inches(0.72), color)
        add_label(sl, name, Inches(0.75), top + Inches(0.1), Inches(2.8), Inches(0.5),
                  size=17, bold=True, color=color)
        add_label(sl, desc, Inches(3.7), top + Inches(0.1), Inches(9.0), Inches(0.5),
                  size=16, color=WHITE)
        top += Inches(0.85)

    add_label(sl,
              "Each layer has one responsibility. The engine knows nothing about queues; nodes know nothing about the heap.",
              Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
              size=14, color=MUTED, italic=True)


# ── Slide 3 — DES Engine ─────────────────────────────────────────────────────

def slide_engine(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, ACCENT)

    add_label(sl, "DES Engine", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
              size=36, bold=True, color=WHITE)
    divider(sl, Inches(1.15))

    # left col — explanation bullets
    bullets = [
        ("Event Calendar", "Priority queue (heapq) sorted by (time, seq)"),
        ("Clock",          "Jumps to next event time — no idle spinning"),
        ("Tiebreaker",     "seq counter ensures deterministic FIFO ordering"),
        ("Dispatch",       "nodes[event.target_id].handle(event)"),
        ("Causality",      "Events schedule future events via clock + delay"),
    ]
    top = Inches(1.45)
    for title, body in bullets:
        add_label(sl, title, Inches(0.6), top, Inches(3.0), Inches(0.32),
                  size=15, bold=True, color=ACCENT)
        add_label(sl, body, Inches(0.6), top + Inches(0.28), Inches(5.6), Inches(0.35),
                  size=14, color=WHITE)
        top += Inches(0.75)

    # right col — event dataclass
    code = (
        "@dataclass(order=True, frozen=True)\n"
        "class Event:\n"
        "    time:      float      # heap key\n"
        "    seq:       int        # tiebreaker\n"
        "    type:      EventType  # compare=False\n"
        "    target_id: str        # compare=False\n"
        "    payload:   Any        # compare=False\n"
        "\n"
        "# Main loop\n"
        "while not scheduler.is_empty():\n"
        "    event      = scheduler.pop_next()\n"
        "    clock      = event.time\n"
        "    nodes[event.target_id].handle(event)"
    )
    add_code_block(sl, code, Inches(6.5), Inches(1.35), Inches(6.5), Inches(4.2), size=13)


# ── Slide 4 — Queueing Networks ───────────────────────────────────────────────

def slide_queueing(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, ACCENT)

    add_label(sl, "Queueing Networks", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
              size=36, bold=True, color=WHITE)
    divider(sl, Inches(1.15))

    # theory table
    headers = ["Model", "ρ = λ/(cμ)", "W (sojourn)", "Wq (wait)"]
    rows = [
        ["M/M/1  (ρ=0.8)", "0.80", "5.00 s", "4.00 s"],
        ["M/M/2  (ρ=0.4)", "0.40", "1.25 s", "0.25 s"],
        ["M/M/4  (ρ=0.2)", "0.20", "1.03 s", "0.03 s"],
    ]
    col_x = [Inches(0.5), Inches(3.5), Inches(6.5), Inches(9.5)]
    col_w = [Inches(2.8), Inches(2.8), Inches(2.8), Inches(2.8)]

    # header row
    for i, h in enumerate(headers):
        add_rect(sl, col_x[i], Inches(1.4), col_w[i], Inches(0.38), RGBColor(0x1E,0x3A,0x5F))
        add_label(sl, h, col_x[i]+Inches(0.1), Inches(1.43), col_w[i], Inches(0.35),
                  size=14, bold=True, color=ACCENT)

    for r, row in enumerate(rows):
        bg = CODE_BG if r % 2 == 0 else RGBColor(0x17,0x23,0x38)
        top = Inches(1.78) + r * Inches(0.4)
        for i, cell in enumerate(row):
            add_rect(sl, col_x[i], top, col_w[i], Inches(0.38), bg)
            add_label(sl, cell, col_x[i]+Inches(0.1), top+Inches(0.05), col_w[i], Inches(0.33),
                      size=13, color=WHITE)

    # code snippet — starts at ~3.0", height 3.9" → bottom at ~6.9" (within 7.5")
    code = (
        "net = QueueingNetwork(warm_up_time=500)\n"
        "# arrival_rate=λ — inter-arrival ~ Exp(λ) by default\n"
        "# override: inter_arrival_fn=lambda: random.uniform(1, 3)\n"
        "net.add_source(\"src\",   arrival_rate=0.8, next_node_id=\"s1\")\n"
        "net.add_server(\"s1\",   service_rate=1.0, c=1)\n"
        "net.add_server(\"s2\",   service_rate=0.8, c=2)\n"
        "net.add_sink(\"sink\")\n"
        "net.add_edge(\"src\", \"s1\")  # weight=1.0 default\n"
        "net.add_edge(\"s1\",  \"s2\")\n"
        "net.add_edge(\"s2\",  \"sink\")\n"
        "net.run(until=100_000, cli=True)"
    )
    add_code_block(sl, code, Inches(0.5), Inches(3.05), Inches(12.3), Inches(3.85), size=13)


# ── Slide 5 — CLI Dashboard ───────────────────────────────────────────────────

def slide_cli(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, ACCENT)

    add_label(sl, "Live CLI Dashboard", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
              size=36, bold=True, color=WHITE)
    divider(sl, Inches(1.15))

    # mock terminal panel
    terminal_bg = RGBColor(0x0A, 0x0E, 0x1A)
    add_rect(sl, Inches(0.5), Inches(1.35), Inches(8.5), Inches(6.3), terminal_bg)
    add_rect(sl, Inches(0.5), Inches(1.35), Inches(8.5), Inches(0.3), RGBColor(0x2D,0x3A,0x4A))

    terminal_lines = [
        ("╭── Network Config ───────────────────────────╮", RGBColor(0x06,0xB6,0xD4)),
        ("│ Node    Kind    Config                       │", RGBColor(0x06,0xB6,0xD4)),
        ("│ source  source  arrivals: Exp(λ=0.5)         │", WHITE),
        ("│ server1 server  c=1  service: Exp(μ=1.0)     │", WHITE),
        ("│ server2 server  c=2  service: Exp(μ=0.8)     │", WHITE),
        ("│ sink    sink    —                            │", WHITE),
        ("╰─────────────────────────────────────────────╯", RGBColor(0x06,0xB6,0xD4)),
        ("", WHITE),
        ("╭── Network Topology ─────────────────────────╮", WHITE),
        ("│  ◉ source  ───────▶  ▣ server1              │", WHITE),
        ("│  ▣ server1 ───────▶  ▣ server2              │", WHITE),
        ("│  ▣ server2 ───────▶  ◎ sink                 │", WHITE),
        ("╰─────────────────────────────────────────────╯", WHITE),
        ("", WHITE),
        ("  t = 42381.50  [████████████░░░░░░░]  42.4%  ", MUTED),
        ("", WHITE),
        ("╭── Live Statistics ──────────────────────────╮", WHITE),
        ("│ Node    Kind   Servers  Queue  ρ     W      │", ACCENT),
        ("│ source  source  λ=0.5    —     —     —      │", WHITE),
        ("│ server1 server  1/1      3   0.50  2.14     │", RGBColor(0x34,0xD3,0x99)),
        ("│ server2 server  2/2      7   0.63  5.61     │", YELLOW),
        ("│ sink    sink     —       —     —  41203done │", GREEN),
        ("╰─────────────────────────────────────────────╯", WHITE),
    ]

    top = Inches(1.72)
    for line, color in terminal_lines:
        txb = sl.shapes.add_textbox(Inches(0.65), top, Inches(8.1), Inches(0.3))
        tf  = txb.text_frame
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(12)
        run.font.color.rgb = color
        run.font.name  = "Courier New"
        top += Inches(0.325)

    # right col — feature callouts
    features = [
        (RGBColor(0x06,0xB6,0xD4), "Config Panel",    "Printed once at startup: arrival\ndistribution, service dist, policy"),
        (ACCENT,  "Topology Panel",   "ASCII graph — nodes and routing\nprobabilities rendered at start"),
        (GREEN,   "Progress Bar",     "Simulated time vs end time\nwith percentage complete"),
        (YELLOW,  "Live Stats Table", "Per-node: queue depth, utilization\nρ (color-coded), W, Wq"),
    ]

    top = Inches(1.4)
    for color, title, body in features:
        add_rect(sl, Inches(9.25), top + Inches(0.1), Inches(0.06), Inches(0.9), color)
        add_label(sl, title, Inches(9.45), top + Inches(0.05), Inches(3.6), Inches(0.35),
                  size=15, bold=True, color=color)
        add_label(sl, body, Inches(9.45), top + Inches(0.38), Inches(3.6), Inches(0.6),
                  size=13, color=MUTED)
        top += Inches(1.3)


# ── Slide 6 — Routing ────────────────────────────────────────────────────────

def slide_routing(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, ACCENT)

    add_label(sl, "Routing Policies", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
              size=36, bold=True, color=WHITE)
    divider(sl, Inches(1.15))

    # three policy cards
    cards = [
        (ACCENT,  "ProbabilisticRouter",
                  "Samples next node from a weighted\ndiscrete distribution over outgoing edges.\nDefault policy — sufficient for Jackson networks.",
                  "net.add_edge(\"S1\", \"S2a\", weight=0.4)\nnet.add_edge(\"S1\", \"S2b\", weight=0.6)"),
        (GREEN,   "RoundRobinRouter",
                  "Cycles through successors in fixed order.\nUseful for load balancing across identical\nparallel servers.",
                  "from des.network.routing import RoundRobinRouter\nnet.set_router(RoundRobinRouter())"),
        (YELLOW,  "ClassBasedRouter",
                  "Routes on customer['class'] field.\nEach source tags customers with a class.\nEnables multiclass networks like criss-cross.",
                  "net.set_router(ClassBasedRouter({\n    \"class1\": \"S2\",\n    \"class3\": \"sink3\",\n}))"),
    ]

    for i, (color, title, body, code) in enumerate(cards):
        left = Inches(0.5) + i * Inches(4.25)
        add_rect(sl, left, Inches(1.4), Inches(4.0), Inches(5.7), CODE_BG)
        add_rect(sl, left, Inches(1.4), Inches(4.0), Inches(0.06), color)
        add_label(sl, title, left + Inches(0.15), Inches(1.5), Inches(3.7), Inches(0.45),
                  size=16, bold=True, color=color)
        add_label(sl, body, left + Inches(0.15), Inches(2.05), Inches(3.7), Inches(1.1),
                  size=13, color=MUTED)
        add_code_block(sl, code, left + Inches(0.1), Inches(3.3), Inches(3.8), Inches(1.5), size=12)

    add_label(sl,
              "All routers implement RoutingPolicy.next_node(customer, successors) → str  —  swap without touching any other code.",
              Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
              size=13, color=MUTED, italic=True)


# ── Slide 7 — Criss-Cross Network ────────────────────────────────────────────

def slide_criss_cross(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, ACCENT)

    add_label(sl, "Criss-Cross Network", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
              size=36, bold=True, color=WHITE)
    divider(sl, Inches(1.15))

    add_label(sl, "A multiclass network where two customer classes share a server but follow different routes.",
              Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4),
              size=15, color=MUTED, italic=True)

    # topology diagram (text-based)
    topo_lines = [
        ("class1_src (λ=0.3) ──▶ S1 ──▶ S2 ──▶ sink_class2",  ACCENT),
        ("class3_src (λ=0.2) ──▶ S1 ──▶ sink_class3",         YELLOW),
        ("                        ↑",                           MUTED),
        ("                    shared server",                   MUTED),
    ]
    top = Inches(1.8)
    for line, color in topo_lines:
        add_label(sl, line, Inches(0.8), top, Inches(10), Inches(0.38),
                  size=16, color=color)
        top += Inches(0.38)

    # code block
    code = (
        "net = QueueingNetwork(warm_up_time=500)\n"
        "\n"
        "net.add_source(\"class1_src\", arrival_rate=0.3,\n"
        "               next_node_id=\"S1\", customer_class=\"class1\")\n"
        "net.add_source(\"class3_src\", arrival_rate=0.2,\n"
        "               next_node_id=\"S1\", customer_class=\"class3\")\n"
        "\n"
        "net.add_server(\"S1\", service_rate=1.0, c=1)\n"
        "net.add_server(\"S2\", service_rate=0.8, c=1)\n"
        "net.add_sink(\"sink_class2\")\n"
        "net.add_sink(\"sink_class3\")\n"
        "\n"
        "net.set_router(ClassBasedRouter({\n"
        "    \"class1\": \"S2\",\n"
        "    \"class3\": \"sink_class3\",\n"
        "}))\n"
        "\n"
        "net.run(until=50_000, cli=True)"
    )
    add_code_block(sl, code, Inches(0.5), Inches(3.3), Inches(7.0), Inches(3.9), size=13)

    # callouts right side
    callouts = [
        (ACCENT,  "Two sources",      "Each tags customers with\na class field"),
        (YELLOW,  "Shared server S1", "Both classes compete\nfor the same resource"),
        (GREEN,   "Class-aware exit", "ClassBasedRouter splits\ntraffic at S1 by class"),
    ]
    top = Inches(3.4)
    for color, title, body in callouts:
        add_rect(sl, Inches(7.8), top + Inches(0.1), Inches(0.06), Inches(0.9), color)
        add_label(sl, title, Inches(8.0), top + Inches(0.05), Inches(4.9), Inches(0.35),
                  size=15, bold=True, color=color)
        add_label(sl, body, Inches(8.0), top + Inches(0.38), Inches(4.9), Inches(0.55),
                  size=13, color=MUTED)
        top += Inches(1.25)


# ── Build ─────────────────────────────────────────────────────────────────────

prs = new_prs()
slide_title(prs)
slide_architecture(prs)
slide_engine(prs)
slide_queueing(prs)
slide_cli(prs)
slide_routing(prs)
slide_criss_cross(prs)

out = "/Users/fastcheetah/Desktop/DiscreteEventSimulator/DiscreteEventSimulator.pptx"
prs.save(out)
print(f"Saved → {out}")
