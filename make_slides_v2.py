"""Generate a 9-slide PPTX covering v2 extensions: shared-station scheduling and RL gym."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG       = RGBColor(0x0F, 0x17, 0x2A)
ACCENT   = RGBColor(0x38, 0xBD, 0xF8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
YELLOW   = RGBColor(0xFB, 0xBF, 0x24)
RED      = RGBColor(0xF8, 0x71, 0x71)
CODE_BG  = RGBColor(0x1E, 0x29, 0x3B)
PURPLE   = RGBColor(0xC0, 0x84, 0xFC)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_label(slide, text, left, top, width, height,
              size=20, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.color.rgb = color
    run.font.italic  = italic


def add_code_block(slide, code, left, top, width, height, size=13):
    add_rect(slide, left, top, width, height, CODE_BG)
    txb = slide.shapes.add_textbox(
        left + Inches(0.18), top + Inches(0.15),
        width - Inches(0.36), height - Inches(0.3)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = ACCENT
        run.font.name  = "Courier New"


def divider(slide, top, color=ACCENT):
    add_rect(slide, Inches(0.5), top, Inches(12.33), Inches(0.04), color)


def slide_header(slide, title):
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, ACCENT)
    add_label(slide, title, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              size=36, bold=True, color=WHITE)
    divider(slide, Inches(1.15))


# ── Slide 1 — Title ───────────────────────────────────────────────────────────

def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), H, ACCENT)

    add_label(sl, "Shared-Station Scheduling",
              Inches(0.5), Inches(1.6), Inches(12), Inches(1.2),
              size=50, bold=True, color=WHITE)
    add_label(sl, "v2 extensions: upstream control, scheduling policies, and RL gym",
              Inches(0.5), Inches(3.0), Inches(10), Inches(0.6),
              size=22, color=ACCENT)

    divider(sl, Inches(3.85))

    pills = [
        ("Buffer + Station",      Inches(0.5)),
        ("Scheduling Policies",   Inches(4.3)),
        ("RL Gym Environment",    Inches(8.1)),
    ]
    for label, left in pills:
        add_rect(sl, left, Inches(4.15), Inches(3.4), Inches(0.55), CODE_BG)
        add_label(sl, label, left + Inches(0.15), Inches(4.17), Inches(3.2), Inches(0.5),
                  size=18, color=ACCENT, bold=True)

    add_label(sl, "Additive design — legacy routed servers unchanged",
              Inches(0.5), Inches(5.5), Inches(12), Inches(0.4),
              size=14, color=MUTED, italic=True)


# ── Slide 2 — Agenda ─────────────────────────────────────────────────────────

def slide_agenda(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Agenda")

    BUF_COLOR = RGBColor(0x67, 0xE8, 0xF9)
    topics = [
        (YELLOW,   "1", "Design Question",
                   "Upstream vs downstream control — where should the RL policy live?"),
        (BUF_COLOR,"2", "Upstream: Buffer + Station",
                   "New passive-queue / active-server split with topology constraints."),
        (RED,      "3", "Downstream: Routed Server",
                   "Legacy server with implicit queue; router fires on each departure."),
        (GREEN,    "4", "Scheduling Policies",
                   "Four built-in policies + Callback hook for custom / RL logic."),
        (PURPLE,   "5", "RL Gym Environment",
                   "SharedStationSchedulingGym: obs, action, reward, and training loop."),
        (ACCENT,   "6", "Queue Discipline",
                   "FIFO / FBFS / LBFS — downstream-only, stage-based selection."),
        (MUTED,    "7", "Discipline Comparison",
                   "Simulated M/M/1 stats across all three disciplines vs. theory."),
    ]

    top = Inches(1.45)
    for color, num, title, desc in topics:
        add_rect(sl, Inches(0.5), top, Inches(12.3), Inches(0.68), CODE_BG)
        add_rect(sl, Inches(0.5), top, Inches(0.07), Inches(0.68), color)
        add_label(sl, num, Inches(0.7), top + Inches(0.13), Inches(0.4), Inches(0.42),
                  size=18, bold=True, color=color)
        add_label(sl, title, Inches(1.2), top + Inches(0.08), Inches(3.2), Inches(0.32),
                  size=14, bold=True, color=color)
        add_label(sl, desc, Inches(4.5), top + Inches(0.08), Inches(8.1), Inches(0.52),
                  size=13, color=WHITE)
        top += Inches(0.80)

    add_label(sl,
              "All extensions are additive — legacy routed servers and v1 networks remain unchanged.",
              Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
              size=13, color=MUTED, italic=True)


# ── Slide 3 — Design Question ─────────────────────────────────────────────────

def slide_design_question(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Design Question: Upstream vs Downstream Control")

    # Two-column layout
    col_w = Inches(5.8)

    # Downstream column
    add_rect(sl, Inches(0.5), Inches(1.35), col_w, Inches(5.7), CODE_BG)
    add_rect(sl, Inches(0.5), Inches(1.35), col_w, Inches(0.06), YELLOW)
    add_label(sl, "Downstream Control", Inches(0.65), Inches(1.45), col_w - Inches(0.3), Inches(0.45),
              size=18, bold=True, color=YELLOW)
    add_label(sl, "Decision happens after service finishes",
              Inches(0.65), Inches(1.98), col_w - Inches(0.3), Inches(0.35),
              size=14, color=MUTED, italic=True)

    ds_bullets = [
        (YELLOW,  "Action",   "Where does the completed job go next?"),
        (YELLOW,  "Strength", "Simple semantics, standard event-driven RL loop"),
        (YELLOW,  "Weakness", "Hides queue ownership;\nless faithful for shared-resource scheduling"),
        (PURPLE,  "PPO",      "P(successor | state);\naction = next hop after service"),
    ]
    top = Inches(2.45)
    for lcolor, label, body in ds_bullets:
        add_label(sl, label, Inches(0.7), top, Inches(1.4), Inches(0.34),
                  size=13, bold=True, color=lcolor)
        add_label(sl, body, Inches(2.2), top, col_w - Inches(1.9), Inches(0.52),
                  size=13, color=WHITE)
        top += Inches(0.62)

    # downstream algorithms section
    add_rect(sl, Inches(0.65), top + Inches(0.08), col_w - Inches(0.3), Inches(0.03), RGBColor(0x2D,0x3A,0x4A))
    add_label(sl, "Algorithms that use this surface:",
              Inches(0.7), top + Inches(0.18), col_w - Inches(0.4), Inches(0.32),
              size=12, bold=True, color=YELLOW)
    ds_algos = [
        "JSQ (Join Shortest Queue) — route to least-loaded server",
        "Probabilistic / class-based routing — Jackson networks",
        "ε-greedy / softmax routing — tabular RL baselines",
        "JFCFS load balancing — parallel server dispatching",
    ]
    top += Inches(0.55)
    for algo in ds_algos:
        add_label(sl, f"• {algo}", Inches(0.8), top, col_w - Inches(0.5), Inches(0.3),
                  size=11, color=MUTED)
        top += Inches(0.3)

    # Upstream column
    left2 = Inches(7.03)
    add_rect(sl, left2, Inches(1.35), col_w, Inches(5.7), CODE_BG)
    add_rect(sl, left2, Inches(1.35), col_w, Inches(0.06), GREEN)
    add_label(sl, "Upstream Control", left2 + Inches(0.15), Inches(1.45), col_w - Inches(0.3), Inches(0.45),
              size=18, bold=True, color=GREEN)
    add_label(sl, "Decision happens when a resource becomes idle",
              left2 + Inches(0.15), Inches(1.98), col_w - Inches(0.3), Inches(0.35),
              size=14, color=MUTED, italic=True)

    us_bullets = [
        (GREEN,   "Action",   "Which waiting buffer gets service now?"),
        (GREEN,   "Strength", "Natural for shared-resource and\nre-entrant queueing systems"),
        (GREEN,   "Weakness", "Richer state; post-service branching\nneeds separate mechanism"),
        (PURPLE,  "PPO",      "P(buffer | queue_lens, busy);\nfits SharedStationGym obs/action"),
    ]
    top = Inches(2.45)
    for lcolor, label, body in us_bullets:
        add_label(sl, label, left2 + Inches(0.2), top, Inches(1.4), Inches(0.34),
                  size=13, bold=True, color=lcolor)
        add_label(sl, body, left2 + Inches(1.7), top, col_w - Inches(1.9), Inches(0.52),
                  size=13, color=WHITE)
        top += Inches(0.62)

    # upstream algorithms section
    add_rect(sl, left2 + Inches(0.15), top + Inches(0.08), col_w - Inches(0.3), Inches(0.03), RGBColor(0x2D,0x3A,0x4A))
    add_label(sl, "Algorithms that use this surface:",
              left2 + Inches(0.2), top + Inches(0.18), col_w - Inches(0.4), Inches(0.32),
              size=12, bold=True, color=GREEN)
    us_algos = [
        "SPT / SRPT — shortest (remaining) processing time first",
        "MaxWeight — serve queue with highest weight × length",
        "Whittle / Gittins index — restless bandit index policies",
        "BackPressure — pressure-differential scheduling",
    ]
    top += Inches(0.55)
    for algo in us_algos:
        add_label(sl, f"• {algo}", left2 + Inches(0.3), top, col_w - Inches(0.5), Inches(0.3),
                  size=11, color=MUTED)
        top += Inches(0.3)

    add_label(sl,
              "Both surfaces coexist — additive design means neither abstraction is forced to cover the other's use case.",
              Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
              size=13, color=MUTED, italic=True)


# ── Slide 3 — Buffer + Station Architecture ───────────────────────────────────

def slide_buffer_station(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "New Architecture: Buffer + Station")

    BUF_COLOR = RGBColor(0x67, 0xE8, 0xF9)

    # ── local drawing helpers ─────────────────────────────────────────────────
    T = Inches(0.032)

    def h_line(x1, yc, x2, color):
        add_rect(sl, x1, yc - T // 2, x2 - x1, T, color)

    def v_seg(xc, y1, y2, color):
        add_rect(sl, xc - T // 2, y1, T, y2 - y1, color)

    def arrowhead(x, yc, color):
        add_label(sl, "▶", x - Inches(0.22), yc - Inches(0.15),
                  Inches(0.25), Inches(0.32), size=11, color=color)

    def node_box(left, top, w, h, title, subtitle, accent):
        add_rect(sl, left, top, w, h, CODE_BG)
        add_rect(sl, left, top, w, Inches(0.05), accent)
        add_label(sl, title, left + Inches(0.1), top + Inches(0.07),
                  w - Inches(0.15), Inches(0.28), size=12, bold=True, color=accent)
        if subtitle:
            add_label(sl, subtitle, left + Inches(0.1), top + Inches(0.3),
                      w - Inches(0.15), Inches(0.21), size=10, color=MUTED)

    # ── geometry ──────────────────────────────────────────────────────────────
    x_src, bw_src = Inches(0.7),  Inches(1.3)
    x_buf, bw_buf = Inches(2.35), Inches(1.2)
    x_sta, bw_sta = Inches(4.1),  Inches(1.55)
    x_snk, bw_snk = Inches(5.9),  Inches(0.8)
    bh     = Inches(0.52)
    bh_sta = Inches(1.0)

    y_fast  = Inches(2.0)
    y_slow  = Inches(4.1)
    cy_fast = y_fast + bh // 2
    cy_slow = y_slow + bh // 2
    cy_sta  = (cy_fast + cy_slow) // 2
    y_sta   = cy_sta - bh_sta // 2
    y_snk   = cy_sta - bh // 2

    merge_x = Inches(3.85)

    # ── node boxes ────────────────────────────────────────────────────────────
    node_box(x_src, y_fast, bw_src, bh, "src_fast", "λ=0.9  Exp arrivals", ACCENT)
    node_box(x_src, y_slow, bw_src, bh, "src_slow", "λ=0.4  Exp arrivals", ACCENT)
    node_box(x_buf, y_fast, bw_buf, bh, "B1", "buffer · FIFO", BUF_COLOR)
    node_box(x_buf, y_slow, bw_buf, bh, "B2", "buffer · FIFO", BUF_COLOR)

    add_rect(sl, x_sta, y_sta, bw_sta, bh_sta, CODE_BG)
    add_rect(sl, x_sta, y_sta, bw_sta, Inches(0.05), GREEN)
    add_label(sl, "Station", x_sta + Inches(0.1), y_sta + Inches(0.07),
              bw_sta - Inches(0.15), Inches(0.28), size=12, bold=True, color=GREEN)
    add_label(sl, "μ=1.5  c=1 server", x_sta + Inches(0.1), y_sta + Inches(0.3),
              bw_sta - Inches(0.15), Inches(0.21), size=10, color=MUTED)
    add_label(sl, "pulls from buffers", x_sta + Inches(0.1), y_sta + Inches(0.52),
              bw_sta - Inches(0.15), Inches(0.21), size=9, color=GREEN, italic=True)

    node_box(x_snk, y_snk, bw_snk, bh, "sink", "", MUTED)

    # ── arrows ────────────────────────────────────────────────────────────────
    # sources → buffers  (push, ACCENT)
    h_line(x_src + bw_src, cy_fast, x_buf, ACCENT)
    arrowhead(x_buf, cy_fast, ACCENT)
    h_line(x_src + bw_src, cy_slow, x_buf, ACCENT)
    arrowhead(x_buf, cy_slow, ACCENT)

    # buffers → vertical merge bar  (pull, GREEN)
    h_line(x_buf + bw_buf, cy_fast, merge_x, GREEN)
    h_line(x_buf + bw_buf, cy_slow, merge_x, GREEN)
    v_seg(merge_x, cy_fast, cy_slow, GREEN)

    # merge → station
    h_line(merge_x, cy_sta, x_sta, GREEN)
    arrowhead(x_sta, cy_sta, GREEN)

    # station → sink
    h_line(x_sta + bw_sta, cy_sta, x_snk, MUTED)
    arrowhead(x_snk, cy_sta, MUTED)

    # ── flow annotations ──────────────────────────────────────────────────────
    # "push" placed above the diagram row, clear of all node boxes
    add_label(sl, "push", x_src + bw_src + Inches(0.05), Inches(1.75),
              Inches(0.65), Inches(0.28), size=10, color=ACCENT, italic=True)

    # "pull ←" in the empty vertical band between the two buffer rows,
    # just above the merge→station arrow (cy_sta) and left of the station box
    add_label(sl, "pull ←", Inches(3.60), cy_sta - Inches(0.22),
              Inches(0.50), Inches(0.28), size=10, color=GREEN, italic=True)

    # scheduling-decision label placed below diagram, above column labels
    add_label(sl, "↑ scheduling decision", merge_x - Inches(0.5), Inches(4.72),
              Inches(2.3), Inches(0.28), size=9, color=YELLOW, italic=True)

    # column labels
    for lbl, cx, color in [
        ("Source",  x_src + bw_src // 2, ACCENT),
        ("Buffer",  x_buf + bw_buf // 2, BUF_COLOR),
        ("Station", x_sta + bw_sta // 2, GREEN),
        ("Sink",    x_snk + bw_snk // 2, MUTED),
    ]:
        add_label(sl, lbl, cx - Inches(0.6), Inches(5.5), Inches(1.2), Inches(0.32),
                  size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

    # ── right column: component rows + constraints ─────────────────────────────
    rx, rw = Inches(7.0), Inches(5.8)

    comp_rows = [
        (BUF_COLOR, "Buffer",    "Passive queue. Stores jobs, tracks\nstats, notifies station on arrival."),
        (GREEN,     "Station",   "Active server. Idles until policy fires;\npulls from chosen buffer (c parallel)."),
        (YELLOW,    "Scheduler", "Policy: choose_buffer(station, buffers)\n→ id. 4 built-ins + Callback."),
        (MUTED,     "SCHED evt", "Fired on-demand when station idles —\navoids cascading event storms."),
    ]
    top = Inches(1.4)
    for color, name, desc in comp_rows:
        add_rect(sl, rx, top, rw, Inches(0.72), CODE_BG)
        add_rect(sl, rx, top, Inches(0.07), Inches(0.72), color)
        add_label(sl, name, rx + Inches(0.22), top + Inches(0.1), Inches(1.4), Inches(0.55),
                  size=14, bold=True, color=color)
        add_label(sl, desc, rx + Inches(1.75), top + Inches(0.1), rw - Inches(1.95), Inches(0.55),
                  size=12, color=WHITE)
        top += Inches(0.83)

    add_label(sl, "Topology constraints — validated at build time:",
              rx, top + Inches(0.1), rw, Inches(0.32),
              size=12, bold=True, color=ACCENT)
    top += Inches(0.48)
    for c in [
        "Buffer outgoing edge → Station only",
        "Station incoming edges ← Buffers only",
        "Station must have exactly 1 outgoing edge",
    ]:
        add_label(sl, f"• {c}", rx + Inches(0.15), top, rw - Inches(0.25), Inches(0.32),
                  size=12, color=WHITE)
        top += Inches(0.36)


# ── Slide 5 — Downstream Architecture ───────────────────────────────────────

def slide_downstream_architecture(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Downstream Architecture: Routed Server")

    SRV_COLOR = YELLOW

    T = Inches(0.032)

    def h_line(x1, yc, x2, color):
        add_rect(sl, x1, yc - T // 2, x2 - x1, T, color)

    def v_seg(xc, y1, y2, color):
        add_rect(sl, xc - T // 2, y1, T, y2 - y1, color)

    def arrowhead(x, yc, color):
        add_label(sl, "▶", x - Inches(0.22), yc - Inches(0.15),
                  Inches(0.25), Inches(0.32), size=11, color=color)

    def node_box(left, top, w, h, title, subtitle, accent):
        add_rect(sl, left, top, w, h, CODE_BG)
        add_rect(sl, left, top, w, Inches(0.05), accent)
        add_label(sl, title, left + Inches(0.1), top + Inches(0.07),
                  w - Inches(0.15), Inches(0.28), size=12, bold=True, color=accent)
        if subtitle:
            add_label(sl, subtitle, left + Inches(0.1), top + Inches(0.3),
                      w - Inches(0.15), Inches(0.21), size=10, color=MUTED)

    # ── geometry ──────────────────────────────────────────────────────────────
    x_src,  bw_src  = Inches(0.55), Inches(1.3)
    x_srv,  bw_srv  = Inches(2.25), Inches(1.9)
    x_fork          = Inches(4.35)
    x_snk,  bw_snk  = Inches(5.5),  Inches(0.85)
    bh              = Inches(0.52)

    y_snk_a = Inches(2.0)
    y_snk_b = Inches(4.0)
    cy_a    = y_snk_a + bh // 2
    cy_b    = y_snk_b + bh // 2
    cy_srv  = (cy_a + cy_b) // 2
    y_srv   = cy_srv - Inches(0.5)

    # ── node boxes ────────────────────────────────────────────────────────────
    node_box(x_src, y_srv - Inches(0.0), bw_src, bh, "src", "λ=0.9  Exp arrivals", ACCENT)

    # Server box — taller to show internal queue note
    add_rect(sl, x_srv, y_srv - Inches(0.22), bw_srv, Inches(0.95), CODE_BG)
    add_rect(sl, x_srv, y_srv - Inches(0.22), bw_srv, Inches(0.05), SRV_COLOR)
    add_label(sl, "Server", x_srv + Inches(0.1), y_srv - Inches(0.15),
              bw_srv - Inches(0.15), Inches(0.28), size=12, bold=True, color=SRV_COLOR)
    add_label(sl, "μ=1.0  c=1  MMcServer", x_srv + Inches(0.1), y_srv + Inches(0.08),
              bw_srv - Inches(0.15), Inches(0.21), size=10, color=MUTED)
    add_label(sl, "[implicit queue inside]", x_srv + Inches(0.1), y_srv + Inches(0.28),
              bw_srv - Inches(0.15), Inches(0.21), size=9, color=SRV_COLOR, italic=True)

    node_box(x_snk, y_snk_a, bw_snk, bh, "sink_A", "", MUTED)
    node_box(x_snk, y_snk_b, bw_snk, bh, "sink_B", "", MUTED)

    # ── arrows ────────────────────────────────────────────────────────────────
    # source → server (push)
    h_line(x_src + bw_src, cy_srv, x_srv, ACCENT)
    arrowhead(x_srv, cy_srv, ACCENT)

    # server → fork (departure)
    h_line(x_srv + bw_srv, cy_srv, x_fork, SRV_COLOR)

    # fork → sink_A and sink_B
    v_seg(x_fork, cy_a, cy_b, SRV_COLOR)
    h_line(x_fork, cy_a, x_snk, SRV_COLOR)
    arrowhead(x_snk, cy_a, SRV_COLOR)
    h_line(x_fork, cy_b, x_snk, SRV_COLOR)
    arrowhead(x_snk, cy_b, SRV_COLOR)

    # ── flow annotations ──────────────────────────────────────────────────────
    add_label(sl, "push", x_src + bw_src + Inches(0.05), Inches(1.75),
              Inches(0.65), Inches(0.28), size=10, color=ACCENT, italic=True)

    # "routing decision" label at the fork, below diagram
    add_label(sl, "↑ routing decision", x_fork - Inches(0.45), Inches(4.72),
              Inches(2.0), Inches(0.28), size=9, color=SRV_COLOR, italic=True)

    # column labels
    for lbl, cx, color in [
        ("Source", x_src  + bw_src  // 2, ACCENT),
        ("Server", x_srv  + bw_srv  // 2, SRV_COLOR),
        ("Sinks",  x_snk  + bw_snk  // 2, MUTED),
    ]:
        add_label(sl, lbl, cx - Inches(0.6), Inches(5.5), Inches(1.2), Inches(0.32),
                  size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

    # ── contrast banner ───────────────────────────────────────────────────────
    add_rect(sl, Inches(0.5), Inches(5.9), Inches(6.1), Inches(0.75), RGBColor(0x1E, 0x29, 0x3B))
    add_rect(sl, Inches(0.5), Inches(5.9), Inches(0.07), Inches(0.75), SRV_COLOR)
    add_label(sl,
              "No explicit Buffer nodes. Jobs queue inside the Server.\n"
              "Router fires on each DEPARTURE — decision is where the job goes next.",
              Inches(0.72), Inches(5.97), Inches(5.7), Inches(0.62),
              size=12, color=WHITE)

    # ── right column: component rows ──────────────────────────────────────────
    rx, rw = Inches(7.0), Inches(5.8)

    comp_rows = [
        (ACCENT,     "Source",     "Generates arrivals. Pushes directly\nto the server — no buffer hop."),
        (SRV_COLOR,  "Server",     "Active. Queues jobs internally;\nroutes each departure via policy."),
        (RED,        "Router",     "Policy: route(customer, successors)\n→ next_node_id. Fires on DEPARTURE."),
        (MUTED,      "DEPART evt", "Steps the RL gym loop; observation\nis queue_length + busy per server."),
    ]
    top = Inches(1.4)
    for color, name, desc in comp_rows:
        add_rect(sl, rx, top, rw, Inches(0.72), CODE_BG)
        add_rect(sl, rx, top, Inches(0.07), Inches(0.72), color)
        add_label(sl, name, rx + Inches(0.22), top + Inches(0.1), Inches(1.4), Inches(0.55),
                  size=14, bold=True, color=color)
        add_label(sl, desc, rx + Inches(1.75), top + Inches(0.1), rw - Inches(1.95), Inches(0.55),
                  size=12, color=WHITE)
        top += Inches(0.83)

    # contrast with upstream
    add_label(sl, "Contrast with upstream (Buffer + Station):",
              rx, top + Inches(0.1), rw, Inches(0.32),
              size=12, bold=True, color=ACCENT)
    top += Inches(0.48)
    contrasts = [
        "Upstream: station pulls from explicit Buffer nodes",
        "Downstream: server pushes to successor on departure",
        "Upstream action → which queue gets served next",
        "Downstream action → where completed job goes next",
    ]
    for c in contrasts:
        add_label(sl, f"• {c}", rx + Inches(0.15), top, rw - Inches(0.25), Inches(0.32),
                  size=12, color=WHITE)
        top += Inches(0.36)


# ── Slide 4 — Scheduling Policies ────────────────────────────────────────────

def slide_scheduling_policies(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Scheduling Policies")

    cards = [
        (ACCENT,  "RoundRobinSchedulingPolicy",
                  "Cycles through upstream buffers\nin fixed order. Stateful cursor.",
                  "net.set_station_scheduler(\n    \"S\",\n    RoundRobinSchedulingPolicy()\n)"),
        (GREEN,   "LongestQueueSchedulingPolicy",
                  "Serves the buffer with the\nmost waiting jobs. Greedy.",
                  "net.set_station_scheduler(\n    \"S\",\n    LongestQueueSchedulingPolicy()\n)"),
        (YELLOW,  "FirstNonEmptySchedulingPolicy",
                  "Always tries buffer[0] first;\nfalls through if empty. Deterministic.",
                  "net.set_station_scheduler(\n    \"S\",\n    FirstNonEmptySchedulingPolicy()\n)"),
        (RGBColor(0xC0, 0x84, 0xFC), "CallbackSchedulingPolicy",
                  "Custom function receives\n(station, buffers) → buffer_id.\nFull RL integration point.",
                  "net.set_station_scheduler(\n    \"S\",\n    CallbackSchedulingPolicy(\n        lambda s, b: agent(b)\n    )\n)"),
    ]

    for i, (color, title, body, code) in enumerate(cards):
        left = Inches(0.5) + i * Inches(3.2)
        add_rect(sl, left, Inches(1.4), Inches(3.0), Inches(5.7), CODE_BG)
        add_rect(sl, left, Inches(1.4), Inches(3.0), Inches(0.06), color)
        add_label(sl, title, left + Inches(0.12), Inches(1.52), Inches(2.78), Inches(0.5),
                  size=13, bold=True, color=color)
        add_label(sl, body, left + Inches(0.12), Inches(2.1), Inches(2.78), Inches(0.9),
                  size=12, color=MUTED)
        add_code_block(sl, code, left + Inches(0.08), Inches(3.1), Inches(2.86), Inches(1.85), size=11)

    add_label(sl,
              "All policies implement SchedulingPolicy.choose_buffer(station, buffers) → buffer_id  "
              "—  fallthrough logic built into Station if chosen buffer is empty.",
              Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
              size=13, color=MUTED, italic=True)


# ── Slide 5 — RL Gym Environment ─────────────────────────────────────────────

def slide_rl_gym(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "RL Gym Environment: SharedStationSchedulingGym")

    # Left: spec table
    specs = [
        ("Class",        "SharedStationSchedulingGym(gym.Env)"),
        ("Observation",  "[buffer queue lengths...] + [station busy servers...]"),
        ("Action space", "Discrete(num_upstream_buffers)  —  index of buffer to serve"),
        ("Reward",       "−Σ queue_length × Δt  per step  (integrated waiting work)"),
        ("Termination",  "Event queue empty"),
        ("Truncation",   "max_steps reached"),
        ("Step",         "Advances sim to next SCHEDULING_DECISION for controlled station"),
    ]
    top = Inches(1.4)
    for i, (label, val) in enumerate(specs):
        bg = CODE_BG if i % 2 == 0 else RGBColor(0x17, 0x23, 0x38)
        add_rect(sl, Inches(0.5), top, Inches(12.3), Inches(0.42), bg)
        add_label(sl, label, Inches(0.65), top + Inches(0.06), Inches(2.0), Inches(0.32),
                  size=13, bold=True, color=ACCENT)
        add_label(sl, val,   Inches(2.85), top + Inches(0.06), Inches(9.8), Inches(0.32),
                  size=13, color=WHITE)
        top += Inches(0.44)

    # Code snippet
    code = (
        "env = SharedStationSchedulingGym(\n"
        "    config,\n"
        "    control_station_id=\"SharedStation\",\n"
        "    max_steps=500,\n"
        ")\n"
        "obs, _ = env.reset()\n"
        "# obs shape: (num_buffers + num_stations,)\n"
        "# e.g. [B1_len, B2_len, SharedStation_busy]\n"
        "\n"
        "obs, reward, terminated, truncated, info = env.step(action)\n"
        "# reward = -sum(queue_lengths) * delta_t\n"
        "# info[\"episode_stats\"] on done"
    )
    add_code_block(sl, code, Inches(0.5), Inches(4.65), Inches(7.3), Inches(2.65), size=13)

    # Training result callout
    add_rect(sl, Inches(8.1), Inches(4.65), Inches(4.7), Inches(2.65), CODE_BG)
    add_rect(sl, Inches(8.1), Inches(4.65), Inches(4.7), Inches(0.06), GREEN)
    add_label(sl, "REINFORCE training result",
              Inches(8.25), Inches(4.75), Inches(4.4), Inches(0.4),
              size=14, bold=True, color=GREEN)
    result_lines = [
        "Network: src_fast (λ=0.9) → B1 ╮",
        "         src_slow (λ=0.4) → B2 ┤→ Station (μ=1.5) → sink",
        "",
        "By episode ~400 the agent learns",
        "to balance servicing across B1/B2,",
        "minimising integrated queue length.",
    ]
    top = Inches(5.25)
    for line in result_lines:
        add_label(sl, line, Inches(8.25), top, Inches(4.4), Inches(0.32),
                  size=12, color=WHITE if line else WHITE)
        top += Inches(0.33)


# ── Slide 6 — Queue Discipline ────────────────────────────────────────────────

def slide_queue_discipline(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Queue Discipline: FIFO / FBFS / LBFS")

    # Scope banner — downstream only
    add_rect(sl, Inches(0.5), Inches(1.22), Inches(12.3), Inches(0.48), RGBColor(0x1E, 0x29, 0x3B))
    add_rect(sl, Inches(0.5), Inches(1.22), Inches(0.07), Inches(0.48), YELLOW)
    add_label(sl,
              "Downstream only — applies to legacy Server nodes. "
              "Buffers are always FIFO; upstream priority is expressed via the scheduling policy (which buffer to pull from), not within-buffer ordering.",
              Inches(0.72), Inches(1.28), Inches(11.9), Inches(0.38),
              size=13, color=WHITE)

    cards = [
        (ACCENT,  "FIFO",
                  "First In, First Out\n(original behaviour)",
                  "Default. Jobs served in\narrival order. Stable and\npredictable."),
        (GREEN,   "FBFS",
                  "Fewest / Shortest\nBy Stage",
                  "Selects the job with the\nlowest stage field.\nFavours early-stage work."),
        (YELLOW,  "LBFS",
                  "Longest / Most-advanced\nBy Stage",
                  "Selects the job with the\nhighest stage field.\nFavours near-completion work."),
    ]

    for i, (color, title, subtitle, body) in enumerate(cards):
        left = Inches(0.5) + i * Inches(4.15)
        add_rect(sl, left, Inches(1.8), Inches(3.9), Inches(4.2), CODE_BG)
        add_rect(sl, left, Inches(1.8), Inches(3.9), Inches(0.06), color)
        add_label(sl, title, left + Inches(0.15), Inches(1.92), Inches(3.6), Inches(0.45),
                  size=20, bold=True, color=color)
        add_label(sl, subtitle, left + Inches(0.15), Inches(2.45), Inches(3.6), Inches(0.55),
                  size=13, color=MUTED, italic=True)
        add_label(sl, body, left + Inches(0.15), Inches(3.1), Inches(3.6), Inches(0.8),
                  size=13, color=WHITE)

    code = (
        "# Discipline is set on the server node\n"
        "net.add_server(\"S1\", service_rate=1.0, c=2,\n"
        "               discipline=\"FBFS\")\n"
        "\n"
        "# Jobs carry a stage field (int) set by the source:\n"
        "net.add_source(\"src\", arrival_rate=0.5,\n"
        "               next_node_id=\"S1\")\n"
        "# stage defaults to 0; set in customer dict to use FBFS/LBFS"
    )
    add_code_block(sl, code, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.7), size=13)


# ── Slide 8 — Discipline Comparison ──────────────────────────────────────────

def slide_discipline_stats(prs):
    # Results from run_disciplines.py — M/M/1 λ=0.8 μ=1.0 ρ=0.8
    # warm_up=2000, run_until=12000, seed=42, stage∈{0,1,2,3} uniform random
    data = {
        "FIFO": {"W": 5.179, "Wq": 4.178, "L": 4.177, "Lq": 3.371},
        "FBFS": {"W": 5.173, "Wq": 4.171, "L": 4.177, "Lq": 3.371},
        "LBFS": {"W": 5.173, "Wq": 4.171, "L": 4.177, "Lq": 3.371},
        "Theory": {"W": 5.000, "Wq": 4.000, "L": 4.000, "Lq": 3.200},
    }
    col_colors = {"FIFO": ACCENT, "FBFS": GREEN, "LBFS": YELLOW, "Theory": MUTED}

    sl = blank_slide(prs)
    fill_bg(sl)
    slide_header(sl, "Discipline Comparison: Simulated M/M/1 Stats")

    # Setup info
    add_label(sl,
              "M/M/1 · λ=0.8 · μ=1.0 · ρ=0.8 · warm-up 2 000 · run 12 000 · stage ∈ {0,1,2,3} uniform",
              Inches(0.5), Inches(1.22), Inches(12.3), Inches(0.35),
              size=13, color=MUTED, italic=True)

    # Column layout: label column + 4 data columns
    label_w = Inches(1.8)
    col_w   = Inches(2.5)
    col_x = {
        "FIFO":   Inches(2.1),
        "FBFS":   Inches(4.7),
        "LBFS":   Inches(7.3),
        "Theory": Inches(9.9),
    }

    # Header row
    header_top = Inches(1.7)
    add_rect(sl, Inches(0.5), header_top, Inches(12.3), Inches(0.55), CODE_BG)
    add_label(sl, "Metric", Inches(0.65), header_top + Inches(0.1), label_w, Inches(0.38),
              size=14, bold=True, color=MUTED)
    for disc, x in col_x.items():
        add_label(sl, disc, x, header_top + Inches(0.1), col_w, Inches(0.38),
                  size=14, bold=True, color=col_colors[disc], align=PP_ALIGN.CENTER)

    # Metric rows
    metrics = [
        ("W",  "Mean sojourn time",  "theory: 5.000"),
        ("Wq", "Mean wait time",     "theory: 4.000"),
        ("L",  "Mean system length", "theory: 4.000"),
        ("Lq", "Mean queue length",  "theory: 3.200"),
    ]
    row_top = Inches(2.35)
    for idx, (key, label, hint) in enumerate(metrics):
        bg = CODE_BG if idx % 2 == 0 else RGBColor(0x17, 0x23, 0x38)
        add_rect(sl, Inches(0.5), row_top, Inches(12.3), Inches(0.65), bg)
        add_label(sl, key,   Inches(0.65), row_top + Inches(0.05), Inches(0.7), Inches(0.38),
                  size=18, bold=True, color=WHITE)
        add_label(sl, label, Inches(1.4),  row_top + Inches(0.05), Inches(0.8), Inches(0.28),
                  size=11, color=MUTED)
        add_label(sl, hint,  Inches(1.4),  row_top + Inches(0.33), Inches(0.8), Inches(0.28),
                  size=10, color=MUTED, italic=True)
        for disc, x in col_x.items():
            val = data[disc][key]
            add_label(sl, f"{val:.3f}", x, row_top + Inches(0.1), col_w, Inches(0.45),
                      size=20, bold=True, color=col_colors[disc], align=PP_ALIGN.CENTER)
        row_top += Inches(0.72)

    # Observation callout
    note_top = row_top + Inches(0.15)
    add_rect(sl, Inches(0.5), note_top, Inches(12.3), Inches(1.05), RGBColor(0x1E, 0x29, 0x3B))
    add_rect(sl, Inches(0.5), note_top, Inches(0.07), Inches(1.05), ACCENT)
    add_label(sl, "Why are L and Lq identical across disciplines?",
              Inches(0.72), note_top + Inches(0.08), Inches(11.9), Inches(0.32),
              size=13, bold=True, color=ACCENT)
    add_label(sl,
              "In a single-server M/M/1 queue, Little's Law (L = λ·W) still holds regardless of discipline. "
              "L and Lq are time-averaged totals — the same jobs depart, just in a different order. "
              "Discipline affects which individual jobs wait longer, not the aggregate queue length.",
              Inches(0.72), note_top + Inches(0.43), Inches(11.9), Inches(0.58),
              size=12, color=WHITE)


# ── Build ─────────────────────────────────────────────────────────────────────

prs = new_prs()
slide_title(prs)
slide_agenda(prs)
slide_design_question(prs)
slide_buffer_station(prs)
slide_downstream_architecture(prs)
slide_scheduling_policies(prs)
slide_rl_gym(prs)
slide_queue_discipline(prs)
slide_discipline_stats(prs)

out = "/Users/fastcheetah/Desktop/DiscreteEventSimulator/DiscreteEventSimulator_v2.pptx"
prs.save(out)
print(f"Saved → {out}")
